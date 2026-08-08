"""
Generates the SokoLink end-to-end project flow deck.

    content defined below  ──▶ python-pptx ──▶ SokoLink_Project_Flow.pptx

WHY a script rather than hand-built slides: the deck restates the same facts as
docs/PRODUCTION_PLAN.md and the PDFs. Generating it means a plan change is a
one-line edit and a re-run, instead of a silent drift between four documents.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

OUT = Path(
    r"c:\Users\Administrator\OneDrive\Desktop\Portfolio Projects\Project Startup"
    r"\SOKOLINK\docs\SokoLink_Project_Flow.pptx"
)

# ── Palette ──────────────────────────────────────────────────────────────────
INK = RGBColor(0x16, 0x18, 0x1D)
MUTED = RGBColor(0x5F, 0x66, 0x72)
GREEN = RGBColor(0x04, 0x78, 0x57)
GREEN_SOFT = RGBColor(0xEC, 0xFD, 0xF5)
BLUE = RGBColor(0x1D, 0x4E, 0xD8)
BLUE_SOFT = RGBColor(0xEF, 0xF6, 0xFF)
AMBER = RGBColor(0xB4, 0x53, 0x09)
AMBER_SOFT = RGBColor(0xFF, 0xFB, 0xEB)
PURPLE = RGBColor(0x6D, 0x28, 0xD9)
PURPLE_SOFT = RGBColor(0xF5, 0xF3, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SURFACE = RGBColor(0xF7, 0xF9, 0xFC)
LINE = RGBColor(0xDF, 0xE3, 0xEA)

FONT = "Segoe UI"
MONO = "Consolas"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def text(sld, txt, x, y, w, h, size=14, color=INK, bold=False, align=PP_ALIGN.LEFT,
         font=FONT, spacing=1.15, anchor=MSO_ANCHOR.TOP):
    """Adds a text box. Returns the frame so callers can append paragraphs."""
    tb = sld.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = spacing
    r = p.add_run()
    r.text = str(txt)
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.name = font
    return tf


def para(tf, txt, size=13, color=INK, bold=False, space_before=6,
         font=FONT, align=PP_ALIGN.LEFT, spacing=1.15):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.line_spacing = spacing
    r = p.add_run()
    r.text = txt
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.name = font
    return p


def box(s, x, y, w, h, fill=SURFACE, line=LINE, radius=True, line_w=1.0):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    shp.text_frame.text = ""
    return shp


def bar(s, x, y, w, h, fill):
    shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def header(s, eyebrow, title, accent=GREEN):
    bar(s, 0, 0, 13.333, 0.12, accent)
    text(s, eyebrow.upper(), 0.75, 0.45, 11, 0.3, size=10.5, color=accent, bold=True)
    text(s, title, 0.75, 0.78, 11.8, 0.9, size=30, bold=True)


def footer(s, n):
    text(s, "SokoLink", 0.75, 6.95, 3, 0.3, size=9, color=MUTED)
    text(s, str(n), 12.2, 6.95, 0.5, 0.3, size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def chip(s, x, y, w, h, label, body, fill, accent):
    box(s, x, y, w, h, fill=fill, line=None)
    bar(s, x, y, w, 0.055, accent)
    tf = text(s, label, x + 0.25, y + 0.28, w - 0.5, 0.4, size=15, bold=True, color=accent)
    for line_txt in body:
        para(tf, line_txt, size=11.5, color=INK, space_before=5)


# ─────────────────────────────────────────────────────────────────────────────
# 1 — Title
# ─────────────────────────────────────────────────────────────────────────────
s = slide()
bar(s, 0, 0, 13.333, 7.5, GREEN)
bar(s, 0, 0, 0.35, 7.5, RGBColor(0x02, 0x5A, 0x41))
text(s, "PROJECT FLOW · END TO END", 1.3, 2.25, 10, 0.4, size=13,
     color=RGBColor(0xA7, 0xF3, 0xD0), bold=True)
text(s, "SokoLink", 1.3, 2.75, 10, 1.3, size=62, color=WHITE, bold=True)
text(s, "Where your audience becomes your customers.", 1.3, 4.05, 10, 0.6,
     size=21, color=RGBColor(0xD1, 0xFA, 0xE5))
bar(s, 1.3, 4.85, 2.2, 0.045, RGBColor(0x6E, 0xE7, 0xB7))
text(s, "TikTok discovery  ·  WhatsApp conversation  ·  M-Pesa payment", 1.3, 5.15, 10, 0.4,
     size=13.5, color=RGBColor(0xA7, 0xF3, 0xD0))
text(s, "August 2026  ·  Fredrick Muchoya  ·  Kenya", 1.3, 6.4, 8, 0.4,
     size=11, color=RGBColor(0x6E, 0xE7, 0xB7))

# ─────────────────────────────────────────────────────────────────────────────
# 2 — The problem
# ─────────────────────────────────────────────────────────────────────────────
s = slide()
header(s, "The problem", "Three systems that don't talk to each other", AMBER)

chip(s, 0.75, 2.0, 3.85, 2.5, "TikTok", [
    "Where buyers discover products.",
    "", "The catalogue already exists —",
    "it just isn't shoppable."], SURFACE, AMBER)
chip(s, 4.75, 2.0, 3.85, 2.5, "WhatsApp", [
    "Where the conversation happens.",
    "", "Buyers ask \"Price?\" and wait",
    "hours for a reply."], SURFACE, AMBER)
chip(s, 8.75, 2.0, 3.85, 2.5, "M-Pesa", [
    "Where the money moves.",
    "", "Reconciled by scrolling SMS,",
    "by hand, at midnight."], SURFACE, AMBER)

box(s, 0.75, 4.85, 11.85, 1.45, fill=AMBER_SOFT, line=None)
bar(s, 0.75, 4.85, 0.06, 1.45, AMBER)
tf = text(s, "Every sale is manually stitched together", 1.1, 5.05, 11.2, 0.4,
          size=17, bold=True, color=AMBER)
para(tf, "…by a seller who is also filming, sourcing, packing and delivering. "
         "3–4 hours a day disappear into work that produces nothing new.",
     size=13, color=INK, space_before=8)
footer(s, 2)

# ─────────────────────────────────────────────────────────────────────────────
# 3 — The evidence
# ─────────────────────────────────────────────────────────────────────────────
s = slide()
header(s, "The evidence", "Prices are withheld on purpose", AMBER)

text(s, "Verified against 24 real captions from live Kenyan sellers:", 0.75, 1.85, 11, 0.4, size=15, color=MUTED)

for i, (n, label, col) in enumerate([
    ("0 of 24", "captions mention KSh, bob\nor shilling", AMBER),
    ("3 of 24", "contain any price-like\nnumber at all", AMBER),
    ("23 of 24", "are hashtags plus\n\"WhatsApp me\"", AMBER),
]):
    x = 0.75 + i * 4.0
    box(s, x, 2.5, 3.7, 1.75, fill=SURFACE, line=LINE)
    text(s, n, x, 2.75, 3.7, 0.6, size=32, bold=True, color=col, align=PP_ALIGN.CENTER)
    text(s, label, x + 0.25, 3.4, 3.2, 0.7, size=12, color=MUTED, align=PP_ALIGN.CENTER)

box(s, 0.75, 4.55, 11.85, 0.95, fill=RGBColor(0xF1, 0xF4, 0xF9), line=LINE)
text(s, "All sizes available kindly WhatsApp#whatApp07XXXXXXXX  #foryoupage  #onlyforyou  #fyppppp",
     1.05, 4.78, 11.2, 0.5, size=13, color=INK, font=MONO)

box(s, 0.75, 5.75, 11.85, 1.05, fill=GREEN_SOFT, line=None)
bar(s, 0.75, 5.75, 0.06, 1.05, GREEN)
tf = text(s, "This is the insight the whole product rests on.", 1.1, 5.92, 11.2, 0.35,
          size=15, bold=True, color=GREEN)
para(tf, "Sellers omit the price to force a DM. That tactic creates the bottleneck — "
         "and it means the price must be read from the video, never the caption.",
     size=12.5, color=INK, space_before=6)
footer(s, 3)

# ─────────────────────────────────────────────────────────────────────────────
# 4 — The product
# ─────────────────────────────────────────────────────────────────────────────
s = slide()
header(s, "The product", "Three pillars", GREEN)

chip(s, 0.75, 1.95, 3.85, 3.6, "Soko Commerce", [
    "Selling.", "",
    "• Catalog", "• Orders", "• Payments", "• WhatsApp Storefront"], GREEN_SOFT, GREEN)
chip(s, 4.75, 1.95, 3.85, 3.6, "Soko Intel", [
    "Getting seen.", "",
    "• Scripts", "• Captions", "• Market Insights", "• Competitor Intelligence"], PURPLE_SOFT, PURPLE)
chip(s, 8.75, 1.95, 3.85, 3.6, "Soko AI", [
    "The conversation.", "",
    "• Customer Support", "• Sales Agent", "• Follow-ups", "• Order Assistance"], BLUE_SOFT, BLUE)

text(s, "SokoLedger — on-device M-Pesa reconciliation — is deliberately parked as Phase 2.",
     0.75, 5.85, 11.8, 0.4, size=12.5, color=MUTED)
footer(s, 4)

# ─────────────────────────────────────────────────────────────────────────────
# 5 — Seller onboarding flow
# ─────────────────────────────────────────────────────────────────────────────
s = slide()
header(s, "Flow 1 · Seller", "From a TikTok handle to a live shop", GREEN)

steps = [
    ("1", "Connect", "Seller sends their\nTikTok handle"),
    ("2", "Scrape", "Apify pulls videos,\ncovers, metrics"),
    ("3", "Draft", "AI reads product\n+ price from media"),
    ("4", "Review", "Seller confirms\nor corrects"),
    ("5", "Publish", "Shop goes live\n— price required"),
]
for i, (n, title_, body) in enumerate(steps):
    x = 0.75 + i * 2.42
    box(s, x, 2.1, 2.15, 2.05, fill=SURFACE, line=LINE)
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.82), Inches(2.28), Inches(0.5), Inches(0.5))
    circ.fill.solid()
    circ.fill.fore_color.rgb = GREEN
    circ.line.fill.background()
    circ.shadow.inherit = False
    ctf = circ.text_frame
    ctf.text = n
    ctf.paragraphs[0].alignment = PP_ALIGN.CENTER
    ctf.paragraphs[0].runs[0].font.size = Pt(15)
    ctf.paragraphs[0].runs[0].font.bold = True
    ctf.paragraphs[0].runs[0].font.color.rgb = WHITE
    text(s, title_, x, 2.92, 2.15, 0.35, size=14, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    text(s, body, x + 0.12, 3.3, 1.9, 0.8, size=11, color=MUTED, align=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        text(s, "→", x + 2.13, 2.85, 0.35, 0.4, size=19, color=LINE, align=PP_ALIGN.CENTER)

box(s, 0.75, 4.55, 11.85, 1.85, fill=GREEN_SOFT, line=None)
bar(s, 0.75, 4.55, 0.06, 1.85, GREEN)
tf = text(s, "The guardrail: the agent proposes, code disposes", 1.1, 4.75, 11.2, 0.4,
          size=16, bold=True, color=GREEN)
para(tf, "The AI drafts words and reads a price it can see. It never decides price, stock or payment status.",
     size=13, color=INK, space_before=8)
para(tf, "A misread price sits in a draft the seller reviews — it never reaches a buyer unattended. "
         "Publishing is a deliberate human act that still requires a price.",
     size=13, color=INK, space_before=4)
footer(s, 5)

# ─────────────────────────────────────────────────────────────────────────────
# 6 — Price cascade
# ─────────────────────────────────────────────────────────────────────────────
s = slide()
header(s, "How it reads a price", "A cost-ordered cascade", GREEN)

tiers = [
    ("Tier 1", "Caption + hashtags", "near-free", "Category hints, sometimes sizes.\nRarely a price.", RGBColor(0xE5, 0xE7, 0xEB), MUTED),
    ("Tier 2", "Cover image", "cheap", "Gemini vision reads a price\nprinted on the thumbnail.", GREEN_SOFT, GREEN),
    ("Tier 3", "Video — visual + AUDIO", "expensive", "Gemini watches and listens.\nThe seller says the price in Sheng.", AMBER_SOFT, AMBER),
]
for i, (tier, src, cost, body, fill, accent) in enumerate(tiers):
    y = 1.95 + i * 1.32
    box(s, 0.75, y, 11.85, 1.15, fill=fill, line=None)
    bar(s, 0.75, y, 0.06, 1.15, accent)
    text(s, tier, 1.1, y + 0.13, 1.1, 0.35, size=13, bold=True, color=accent)
    text(s, src, 2.25, y + 0.13, 3.6, 0.35, size=15, bold=True, color=INK)
    text(s, cost, 2.25, y + 0.55, 3.6, 0.35, size=11.5, color=accent, bold=True)
    text(s, body, 6.1, y + 0.16, 6.3, 0.85, size=12, color=INK)

box(s, 0.75, 6.0, 11.85, 0.85, fill=SURFACE, line=LINE)
tf = text(s, "Escalates only when the tier above fails. Every video is processed once ever, "
             "keyed by video id, and cached.", 1.1, 6.15, 11.2, 0.35, size=13, color=INK)
para(tf, "Proven live on real clips: KES 500 / 550 / 550 / 600 read correctly.",
     size=12, color=GREEN, bold=True, space_before=4)
footer(s, 6)

# ─────────────────────────────────────────────────────────────────────────────
# 7 — Buyer journey
# ─────────────────────────────────────────────────────────────────────────────
s = slide()
header(s, "Flow 2 · Buyer", "Buying without leaving WhatsApp", BLUE)

journey = [
    ("WhatsApp chat", "Buyer asks about\na product", BLUE),
    ("Browse products", "Agent answers,\nsends a link", BLUE),
    ("In-app browser", "WhatsApp opens the\nstorefront inline", BLUE),
    ("SokoLink storefront", "Mobile-first,\nserver-rendered", GREEN),
    ("M-Pesa checkout", "STK push to the\nsame phone", GREEN),
    ("Return to WhatsApp", "Order confirmed\nin the thread", GREEN),
]
for i, (title_, body, accent) in enumerate(journey):
    col, row = i % 3, i // 3
    x = 0.75 + col * 4.05
    y = 1.95 + row * 2.15
    box(s, x, y, 3.7, 1.75, fill=SURFACE, line=LINE)
    bar(s, x, y, 3.7, 0.05, accent)
    text(s, f"{i + 1}", x + 0.22, y + 0.22, 0.4, 0.3, size=11, bold=True, color=accent)
    text(s, title_, x + 0.22, y + 0.6, 3.3, 0.35, size=15, bold=True, color=INK)
    text(s, body, x + 0.22, y + 1.0, 3.3, 0.65, size=11.5, color=MUTED)
    if col < 2:
        text(s, "→", x + 3.7, y + 0.65, 0.35, 0.4, size=19, color=LINE, align=PP_ALIGN.CENTER)

box(s, 0.75, 6.25, 11.85, 0.72, fill=BLUE_SOFT, line=None)
bar(s, 0.75, 6.25, 0.06, 0.72, BLUE)
text(s, "Not WhatsApp Flows — the ordinary in-app browser. No Flows approval dependency, "
        "and full control of the interface.", 1.1, 6.42, 11.2, 0.4, size=13, color=INK)
footer(s, 7)

# ─────────────────────────────────────────────────────────────────────────────
# 8 — The three hard parts
# ─────────────────────────────────────────────────────────────────────────────
s = slide()
header(s, "The hard parts", "Three problems to design deliberately", BLUE)

hard = [
    ("Identity across the hop",
     "The storefront must know which buyer and which shop — with no login.",
     "Signed, scoped, short-lived token in the link. Never a phone number in a URL: "
     "links leak via history, referrers and forwarding."),
    ("M-Pesa interrupts the browser",
     "The STK prompt takes over the same phone the webview is running on.",
     "Checkout polls order status and renders paid / failed / timed-out honestly on return. "
     "The callback stays the only payment truth."),
    ("The return hop",
     "A confirmed order should not strand the buyer on a success page.",
     "Deep-link back into the conversation with the seller. Designed in, not bolted on."),
]
for i, (t, problem, fix) in enumerate(hard):
    y = 1.9 + i * 1.5
    box(s, 0.75, y, 11.85, 1.34, fill=SURFACE, line=LINE)
    bar(s, 0.75, y, 0.06, 1.34, AMBER)
    text(s, t, 1.1, y + 0.13, 4.5, 0.35, size=14.5, bold=True, color=INK)
    text(s, problem, 1.1, y + 0.57, 4.4, 0.7, size=11, color=MUTED)
    text(s, fix, 6.0, y + 0.18, 6.4, 1.05, size=11.5, color=INK)

text(s, "These work perfectly on a desktop test and fall apart on a real handset — "
        "which is why every milestone is tested on a real phone.",
     0.75, 6.5, 11.8, 0.4, size=12, color=MUTED)
footer(s, 8)

# ─────────────────────────────────────────────────────────────────────────────
# 9 — Soko Intel / the moat
# ─────────────────────────────────────────────────────────────────────────────
s = slide()
header(s, "Soko Intel", "The moat is the corpus, not the feature", PURPLE)

loop = [
    ("Niche search", "Seller types a\nkeyword"),
    ("Outlier posts", "Which posts broke\nout, and why"),
    ("Extract hooks", "Tagged by niche +\nlanguage register"),
    ("Hook corpus", "Hooks that worked\nhere, in Sheng"),
    ("Generate", "Grounded in the corpus\n+ the seller's stock"),
]
for i, (t, b) in enumerate(loop):
    x = 0.75 + i * 2.42
    box(s, x, 2.05, 2.15, 1.75, fill=PURPLE_SOFT, line=None)
    bar(s, x, 2.05, 2.15, 0.05, PURPLE)
    text(s, t, x + 0.12, 2.3, 1.9, 0.4, size=13, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)
    text(s, b, x + 0.12, 2.78, 1.9, 0.8, size=10.5, color=MUTED, align=PP_ALIGN.CENTER)
    if i < len(loop) - 1:
        text(s, "→", x + 2.13, 2.65, 0.35, 0.4, size=19, color=RGBColor(0xC4, 0xB5, 0xFD), align=PP_ALIGN.CENTER)

box(s, 0.75, 4.15, 5.8, 2.15, fill=SURFACE, line=LINE)
tf = text(s, "Why it defends", 1.1, 4.35, 5.2, 0.35, size=15, bold=True, color=PURPLE)
para(tf, "Hooks are culture- and language-bound. A hook that converts in "
         "Los Angeles does not convert in Nairobi.", size=12, color=INK, space_before=8)
para(tf, "Existing tools are built for other markets. A competitor can copy the "
         "feature — but not the corpus.", size=12, color=INK, space_before=6)

box(s, 6.8, 4.15, 5.8, 2.15, fill=AMBER_SOFT, line=None)
bar(s, 6.8, 4.15, 0.06, 2.15, AMBER)
tf = text(s, "Grounding, not training", 7.15, 4.35, 5.2, 0.35, size=15, bold=True, color=AMBER)
para(tf, "We are not fine-tuning a model. We retrieve real high-performing hooks "
         "into the prompt as examples.", size=12, color=INK, space_before=8)
para(tf, "A fraction of the cost — and it improves the moment a new hook lands, "
         "with no retraining cycle.", size=12, color=INK, space_before=6)
footer(s, 9)

# ─────────────────────────────────────────────────────────────────────────────
# 10 — Tech stack
# ─────────────────────────────────────────────────────────────────────────────
s = slide()
header(s, "Tech stack", "Chosen to be debuggable", BLUE)

rows = [
    ("Language", "Python 3.11+", "The maintainer's strongest language"),
    ("Framework", "FastAPI", "Async, typed, every endpoint visible"),
    ("Validation", "Pydantic v2", "API schemas AND LLM output schemas"),
    ("Database", "PostgreSQL + SQLAlchemy + Alembic", "Constraints as rails; migrations as version control"),
    ("Templating", "Jinja2", "Server-rendered storefront, minimal client JS"),
    ("AI", "Gemini (google-genai)", "Best on Sheng/Swahili — printed and spoken"),
    ("Scraping", "Apify, behind our own adapter", "Engine swap is a one-file change"),
    ("Messaging", "WhatsApp Cloud API", "Where Kenyan commerce already lives"),
    ("Payments", "M-Pesa Daraja (STK push)", "Universal, including the unbanked"),
    ("Hosting", "Railway", "App and database on one platform"),
]
y0 = 1.8
for i, (layer, choice, why) in enumerate(rows):
    y = y0 + i * 0.42
    if i % 2 == 0:
        box(s, 0.75, y, 11.85, 0.40, fill=SURFACE, line=None, radius=False)
    text(s, layer, 0.95, y + 0.05, 1.9, 0.3, size=11, color=MUTED, bold=True)
    text(s, choice, 2.95, y + 0.05, 4.3, 0.3, size=11.5, color=INK, bold=True)
    text(s, why, 7.4, y + 0.05, 5.1, 0.3, size=11, color=MUTED)

box(s, 0.75, 6.15, 11.85, 0.6, fill=BLUE_SOFT, line=None)
bar(s, 0.75, 6.15, 0.06, 0.6, BLUE)
text(s, "No AI framework. No vector database. Provider SDKs called directly — "
        "so nothing is magic, and nothing is un-debuggable.",
     1.1, 6.29, 11.2, 0.35, size=12.5, color=INK)
footer(s, 10)

# ─────────────────────────────────────────────────────────────────────────────
# 11 — Roadmap
# ─────────────────────────────────────────────────────────────────────────────
s = slide()
header(s, "Roadmap", "Four phases, nine milestones", GREEN)

phases = [
    ("Phase 1 · Soko Commerce", GREEN, GREEN_SOFT,
     ["P0  Foundation", "P1  Catalog + draft agent", "P2  WhatsApp channel",
      "P3  Storefront", "P4  Orders + Payments"]),
    ("Phase 2 · Soko AI", BLUE, BLUE_SOFT,
     ["P5  Sales agent + support", "P6  Follow-ups + order assistance"]),
    ("Phase 3 · Soko Intel", PURPLE, PURPLE_SOFT,
     ["P7  Competitor intelligence", "P8  Scripts + captions"]),
    ("Phase 4 · Launch", AMBER, AMBER_SOFT,
     ["P9  Hardening + pilot"]),
]
x = 0.75
for title_, accent, fill, items in phases:
    w = 2.87
    box(s, x, 1.95, w, 3.5, fill=fill, line=None)
    bar(s, x, 1.95, w, 0.06, accent)
    text(s, title_, x + 0.2, 2.18, w - 0.4, 0.6, size=13, bold=True, color=accent)
    tf = None
    for j, it in enumerate(items):
        yy = 2.95 + j * 0.42
        text(s, it, x + 0.2, yy, w - 0.35, 0.35, size=11.5, color=INK)
    del tf
    x += 3.02

box(s, 0.75, 5.7, 11.85, 1.1, fill=SURFACE, line=LINE)
tf = text(s, "Commerce first — nothing matters until a seller can sell.", 1.1, 5.87, 11.2, 0.35,
          size=14, bold=True, color=INK)
para(tf, "Soko Intel is independent of WhatsApp, so it can be built in parallel while "
         "Meta business verification is pending.", size=12.5, color=MUTED, space_before=6)
footer(s, 11)

# ─────────────────────────────────────────────────────────────────────────────
# 12 — How we work
# ─────────────────────────────────────────────────────────────────────────────
s = slide()
header(s, "How we work", "Documented as we build", GREEN)

cards = [
    ("Every file explains itself", GREEN, GREEN_SOFT,
     "A header docstring on every file: what it does, the pipeline it sits in, and "
     "why it is shaped that way. In six weeks nobody remembers why the cascade "
     "escalates in that order — the code has to say so."),
    ("Tests ship with the logic", BLUE, BLUE_SOFT,
     "Never \"tests later\". Money paths get explicit replay tests proving idempotency. "
     "External services are always mocked — tests never hit a paid API."),
    ("Rails before agent", AMBER, AMBER_SOFT,
     "The payment path is built and proven with a plain button before the sales agent "
     "may call it. The agent requests a charge; deterministic code executes it."),
    ("Cost rails are not optional", PURPLE, PURPLE_SOFT,
     "Anything a user can trigger repeatedly gets a cache, a quota and — where "
     "appropriate — a paywall, designed in from the start rather than after a bill."),
]
for i, (t, accent, fill, body) in enumerate(cards):
    col, row = i % 2, i // 2
    x = 0.75 + col * 6.05
    y = 1.95 + row * 2.35
    box(s, x, y, 5.8, 2.1, fill=fill, line=None)
    bar(s, x, y, 0.06, 2.1, accent)
    text(s, t, x + 0.35, y + 0.22, 5.2, 0.4, size=15, bold=True, color=accent)
    text(s, body, x + 0.35, y + 0.72, 5.2, 1.2, size=12, color=INK)
footer(s, 12)

# ─────────────────────────────────────────────────────────────────────────────
# 13 — Next
# ─────────────────────────────────────────────────────────────────────────────
s = slide()
bar(s, 0, 0, 13.333, 7.5, GREEN)
text(s, "WHAT HAPPENS NEXT", 1.3, 1.5, 10, 0.4, size=13,
     color=RGBColor(0xA7, 0xF3, 0xD0), bold=True)
text(s, "Start building", 1.3, 1.95, 10, 1.0, size=46, color=WHITE, bold=True)

nexts = [
    ("Now", "P0 Foundation — FastAPI spine, own database, CI green"),
    ("Now, in parallel", "P7 Apify niche-search spike — the corpus starts here, and it does not wait on Meta"),
    ("On approval", "P2 WhatsApp channel — developer account approved, business verification pending"),
]
for i, (when, what) in enumerate(nexts):
    y = 3.35 + i * 0.95
    box(s, 1.3, y, 10.7, 0.78, fill=RGBColor(0x03, 0x69, 0x4C), line=None)
    text(s, when, 1.6, y + 0.22, 2.4, 0.35, size=12.5, bold=True, color=RGBColor(0x6E, 0xE7, 0xB7))
    text(s, what, 4.2, y + 0.22, 7.5, 0.35, size=13, color=WHITE)

text(s, "Every implementation documented — so the codebase is understandable, "
        "debuggable, and safe to change.",
     1.3, 6.45, 10.7, 0.4, size=12.5, color=RGBColor(0xA7, 0xF3, 0xD0))

prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Slides: {len(prs.slides.__iter__.__self__._sldIdLst)}")
